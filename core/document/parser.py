import csv
import json
import os
import re
from typing import Optional

import yaml


# 支持的文档类型（覆盖企业常见格式）
SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".py", ".js", ".ts", ".java", ".cpp", ".h", ".c",
    ".json", ".yaml", ".yml", ".xml", ".html", ".htm", ".css",
    ".csv", ".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".ppt",
    ".log", ".ini", ".cfg", ".conf", ".toml", ".rs", ".go", ".rb",
    ".php", ".sql", ".sh", ".bat", ".ps1", ".rtf", ".eml",
}

TEXT_EXTENSIONS = {
    ".txt", ".md", ".py", ".js", ".ts", ".java", ".cpp", ".h", ".c",
    ".json", ".yaml", ".yml", ".xml", ".html", ".htm", ".css",
    ".log", ".ini", ".cfg", ".conf", ".toml", ".rs", ".go", ".rb",
    ".php", ".sql", ".sh", ".bat", ".ps1", ".rtf", ".eml",
}

# 文档分类映射
CATEGORY_MAP = {
    ".txt": "text", ".md": "document", ".pdf": "document",
    ".docx": "document", ".xlsx": "spreadsheet", ".xls": "spreadsheet",
    ".csv": "data", ".pptx": "presentation", ".ppt": "presentation",
    ".json": "data", ".yaml": "data", ".yml": "data",
    ".xml": "data", ".html": "web", ".htm": "web",
    ".py": "code", ".js": "code", ".ts": "code",
    ".java": "code", ".cpp": "code", ".c": "code", ".h": "code",
    ".rs": "code", ".go": "code", ".rb": "code", ".php": "code",
    ".sql": "code", ".sh": "code", ".bat": "code", ".ps1": "code",
}


def get_file_type(file_path: str) -> str:
    """获取文件类型（小写扩展名）"""
    _, ext = os.path.splitext(file_path)
    return ext.lower()


def get_doc_category(file_path: str) -> str:
    """获取文档分类"""
    ext = get_file_type(file_path)
    return CATEGORY_MAP.get(ext, "other")


def is_supported(file_path: str) -> bool:
    """检查文件是否支持解析"""
    ext = get_file_type(file_path)
    return ext in SUPPORTED_EXTENSIONS


def is_text_file(file_path: str) -> bool:
    ext = get_file_type(file_path)
    return ext in TEXT_EXTENSIONS


def parse_document(file_path: str, extract_tables: bool = True) -> Optional[str]:
    """解析文档，返回文本内容

    Args:
        file_path: 文件路径
        extract_tables: 是否提取表格（PDF专用）
    """
    ext = get_file_type(file_path)

    if ext in TEXT_EXTENSIONS:
        return _parse_text_file(file_path)
    elif ext == ".pdf":
        return _parse_pdf(file_path, extract_tables=extract_tables)
    elif ext == ".docx":
        return _parse_docx(file_path)
    elif ext == ".xlsx":
        return _parse_xlsx(file_path)
    elif ext == ".xls":
        return _parse_xls(file_path)
    elif ext == ".pptx":
        return _parse_pptx(file_path)
    elif ext == ".ppt":
        return _parse_ppt(file_path)

    return None


def _parse_text_file(file_path: str) -> Optional[str]:
    """解析纯文本文件（自动检测编码）"""
    # 用 chardet 检测编码
    try:
        import chardet
        with open(file_path, "rb") as f:
            raw = f.read(10000)
        detected = chardet.detect(raw)
        encoding = detected.get("encoding", "utf-8")
        # chardet 有时会检测出 ISO-8859-1，用 utf-8 再试
        if encoding and encoding.lower() in ("ascii", "iso-8859-1", "windows-1252"):
            encoding = "utf-8"
    except ImportError:
        encoding = None

    encodings = []
    if encoding:
        encodings.append(encoding)
    encodings.extend(["utf-8", "gbk", "gb2312", "latin-1"])

    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    return None


def _parse_pdf(file_path: str, extract_tables: bool = True) -> Optional[str]:
    """解析 PDF 文件（含表格提取、元数据）"""
    parts = []

    # 提取元数据
    try:
        from pypdf import PdfReader as PdfMetaReader
        reader = PdfMetaReader(file_path)
        meta = reader.metadata
        if meta:
            meta_lines = []
            if meta.title:
                meta_lines.append(f"标题: {meta.title}")
            if meta.author:
                meta_lines.append(f"作者: {meta.author}")
            if meta.subject:
                meta_lines.append(f"主题: {meta.subject}")
            if meta_lines:
                parts.append("=== 文档信息 ===")
                parts.extend(meta_lines)
                parts.append("")
    except ImportError:
        pass
    except Exception:
        pass

    try:
        import pdfplumber

        with pdfplumber.open(file_path) as pdf:
            parts.append(f"=== 共 {len(pdf.pages)} 页 ===")
            for i, page in enumerate(pdf.pages, 1):
                page_text = []
                # 文本
                text = page.extract_text()
                if text and text.strip():
                    page_text.append(text)

                # 表格（结构化输出）
                if extract_tables:
                    tables = page.extract_tables()
                    for j, table in enumerate(tables):
                        if table and len(table) > 0:
                            page_text.append(f"\n--- 表格 {j + 1} ---")
                            for row in table:
                                # 过滤空行
                                clean = [c.strip() if c else "" for c in row]
                                if any(clean):
                                    page_text.append(" | ".join(clean))

                if page_text:
                    parts.append(f"\n--- 第 {i} 页 ---")
                    parts.extend(page_text)

            result = "\n".join(parts)
            return result if result.strip() else None
    except ImportError:
        pass
    except Exception as e:
        parts.append(f"[PDF解析部分失败: {e}]")

    # 回退到 PyPDF2
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        combined = "\n\n".join(pages)
        return combined if combined.strip() else parts and "\n".join(parts) or None
    except ImportError:
        pass
    except Exception:
        pass

    return "\n".join(parts) if any(p.strip() for p in parts) else None


def _parse_docx(file_path: str) -> Optional[str]:
    """解析 Word 文档"""
    try:
        from docx import Document

        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)

        # 也提取表格内容
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                paragraphs.append(" | ".join(cells))

        return "\n".join(paragraphs) if paragraphs else None
    except Exception:
        return None


def _parse_xlsx(file_path: str) -> Optional[str]:
    """解析 Excel 文件"""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True, data_only=True)
        rows = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows.append(f"=== 工作表: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = "\t".join(cells)
                if line.strip():
                    rows.append(line)
        return "\n".join(rows) if rows else None
    except Exception:
        return None


def _parse_xls(file_path: str) -> Optional[str]:
    """解析旧版 Excel (.xls) 文件"""
    try:
        import pandas as pd
        sheets = pd.read_excel(file_path, sheet_name=None)
        rows = []
        for name, df in sheets.items():
            rows.append(f"=== 工作表: {name} ===")
            rows.append(df.to_string(index=False))
        return "\n".join(rows)
    except ImportError:
        return f"[错误: 需要安装 pandas 来解析 .xls 文件]"
    except Exception:
        return None


def _parse_pptx(file_path: str) -> Optional[str]:
    """解析 PowerPoint 文件"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation(file_path)
        parts = [f"=== 共 {len(prs.slides)} 张幻灯片 ==="]

        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"\n--- 幻灯片 {i} ---"]

            # 提取所有文本框内容
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            # 保留缩进层级
                            prefix = "  " * para.level if para.level else ""
                            slide_parts.append(f"{prefix}{text}")

                # 提取表格内容
                if shape.has_table:
                    table = shape.table
                    slide_parts.append("\n[表格]:")
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_parts.append(" | ".join(cells))

                # 提取图表中的文字
                if hasattr(shape, 'has_chart') and shape.has_chart:
                    slide_parts.append("[图表]")

                # 提取备注
                if hasattr(shape, 'has_notes_slide') and shape.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_parts.append(f"备注: {notes}")

            parts.append("\n".join(slide_parts))

        return "\n".join(parts)
    except ImportError:
        return "[错误: 需要安装 python-pptx 来解析 .pptx 文件]"
    except Exception as e:
        return f"[PPTX解析失败: {e}]"


def _parse_ppt(file_path: str) -> Optional[str]:
    """解析旧版 PowerPoint (.ppt) 文件（需先转换）"""
    return f"[提示: 旧版 .ppt 格式需要先转换为 .pptx 才能解析]"


def chunk_text(
    text: str, chunk_size: int = 500, overlap: int = 100
) -> list[str]:
    """将文本分割成块，支持智能分段"""
    if not text:
        return []

    # 先按段落分割
    paragraphs = re.split(r"\n\s*\n", text.strip())
    chunks = []
    current_chunk = []

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        # 如果段落本身就很大，按句子分割
        if len(para) > chunk_size:
            if current_chunk:
                chunks.append("\n\n".join(current_chunk))
                current_chunk = []

            sentences = re.split(r"(?<=[.!?。！？])\s+", para)
            temp = ""
            for sent in sentences:
                if len(temp) + len(sent) > chunk_size and temp:
                    chunks.append(temp)
                    temp = sent
                else:
                    temp = (temp + " " + sent) if temp else sent
            if temp:
                chunks.append(temp)
            continue

        current_len = sum(len(p) for p in current_chunk) + len(para)
        if current_len > chunk_size and current_chunk:
            chunks.append("\n\n".join(current_chunk))
            # 保留最后一小段做 overlap
            overlap_texts = []
            overlap_len = 0
            for p in reversed(current_chunk):
                if overlap_len + len(p) > overlap and overlap_texts:
                    break
                overlap_texts.insert(0, p)
                overlap_len += len(p)
            current_chunk = overlap_texts

        current_chunk.append(para)

    if current_chunk:
        chunks.append("\n\n".join(current_chunk))

    return chunks if chunks else [text]


def get_file_metadata(file_path: str) -> dict:
    """获取文件元数据"""
    try:
        stat = os.stat(file_path)
        _, ext = os.path.splitext(file_path)
        return {
            "size": stat.st_size,
            "modified": stat.st_mtime,
            "created": getattr(stat, "st_ctime", 0),
            "ext": ext.lower(),
            "category": CATEGORY_MAP.get(ext.lower(), "other"),
        }
    except Exception:
        return {"size": 0, "ext": "", "category": "other"}


def count_tokens(text: str) -> int:
    """简单估算 token 数（中英文混合）"""
    if not text:
        return 0
    # 中文字符算 1.5 token，英文按空格分词
    chinese_chars = len(re.findall(r"[一-鿿]", text))
    english_words = len(re.findall(r"[a-zA-Z0-9]+", text))
    return int(chinese_chars * 1.5 + english_words * 1.3)
