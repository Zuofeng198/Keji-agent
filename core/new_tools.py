"""扩展工具集 —— 代码执行、文件浏览、知识库、文档生成、数据处理"""

import os
import sys
import csv
import io
import json
import re
import datetime
import statistics
import tempfile
import subprocess
from typing import Optional

from core.tools import register_tool
from core.document.parser import parse_document, is_supported, get_file_metadata
from core.document.indexer import get_indexer
from core.rag.vector_store import get_vector_store
from core.database.db import get_db


# ═══════════════════════════════════════════════════════════════
# 核心工具：模型直接写 Python 代码执行
# ═══════════════════════════════════════════════════════════════

@register_tool(
    name="run_code",
    description="执行 Python 代码完成任意任务。可导入本项目的工具函数：create_document/create_table/create_presentation/create_folder/delete_file/browse_files/search_files/read_file/read_document/analyze_data/format_data/query_knowledge/index_knowledge",
    parameters={
        "code": {"type": "string", "description": "Python 代码，print() 输出结果"},
    },
    category="utility",
    timeout=600,
)
def run_code(code: str) -> str:
    # 确保 core 模块可导入
    project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))

    tmp = tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False, encoding="utf-8")
    # 分段写入，避免 f-string 花括号转义问题
    tmp.write("import sys\n")
    tmp.write(f"sys.path.insert(0, {project_root!r})\n")
    tmp.write("sys.stdout.reconfigure(encoding='utf-8', errors='replace')\n")
    tmp.write("""
try:
    import openpyxl.cell.cell as _oc
    _orig_cell = openpyxl.Worksheet.cell
    def _safe_cell(self, row, column, value=None):
        c = _orig_cell(self, row, column)
        if value is not None and not isinstance(c, _oc.MergedCell):
            c.value = value
        return c
    openpyxl.Worksheet.cell = _safe_cell
except Exception:
    pass

try:
""")
    # 缩进用户代码，放在 try: 块内
    for line in code.split('\n'):
        tmp.write('    ' + line + '\n')
    tmp.write("""
except AttributeError as _e:
    if 'MergedCell' in str(_e):
        print("\\nMergedCell error: cannot assign to merged cell.")
        print("Use isinstance(cell, openpyxl.cell.cell.MergedCell) to skip.")
        print("Detail: " + str(_e))
    else:
        raise

sys.stdout.flush()
""")
    tmp.close()

    try:
        proc = subprocess.run(
            [sys.executable, tmp.name],
            capture_output=True, text=True, timeout=300,
            encoding="utf-8", errors="replace",
            env={**os.environ, "PYTHONIOENCODING": "utf-8"},
        )
        out = proc.stdout.strip()
        err = proc.stderr.strip()

        # 提取关键错误信息（取最后一行 Error + 最后 4 行堆栈）
        err_text = ""
        if err:
            err_lines = err.split('\n')
            # 找最后一个含 "Error" 的行（如 AttributeError: ...）
            last_error = ""
            for line in reversed(err_lines):
                if 'Error' in line:
                    last_error = line.strip()
                    break
            # 堆栈尾部（文件:行号 + 实际代码）
            tail = [l for l in err_lines[-5:] if l.strip()]
            err_text = (last_error + '\n' + '\n'.join(tail)) if last_error else '\n'.join(tail)
            err_text = err_text[-800:]  # 最多 800 字符

        if not out and err_text:
            out = "代码执行出错：\n" + err_text
        elif not out:
            out = "(无输出 —— 请检查代码中是否有 print() 语句)"
        elif err_text:
            out += "\n\n[执行警告]\n" + err_text

        # 过滤掉 JSON 格式的日志行
        clean_lines = []
        for _line in out.split('\n'):
            _s = _line.strip()
            if _s.startswith('{"timestamp"') and '"level"' in _s and '"logger"' in _s:
                continue
            clean_lines.append(_line)
        out = '\n'.join(clean_lines).strip() or out[:200]
        return out[:4000]
    except subprocess.TimeoutExpired:
        return "错误：代码执行超时（300秒）"
    finally:
        try:
            os.unlink(tmp.name)
        except OSError:
            pass


# ──────────── 文件浏览工具 ────────────


@register_tool(
    name="browse_files",
    description="浏览电脑上的文件夹内容，列出目录下的文件和子文件夹",
    parameters={
        "path": {
            "type": "string",
            "description": "文件夹路径，如 C:\\Users\\ 或 D:\\work\\，默认 Desktop",
        },
    },
    category="filesystem",
    timeout=15,
)
def browse_files(path: str = "") -> str:
    if not path:
        path = os.path.expanduser("~\\Desktop")
    path = os.path.abspath(path)

    if not os.path.exists(path):
        return f"错误：路径不存在「{path}」"
    if not os.path.isdir(path):
        return f"错误：不是文件夹「{path}」"

    try:
        items = []
        for name in sorted(os.listdir(path)):
            full = os.path.join(path, name)
            try:
                stat = os.stat(full)
                is_dir = os.path.isdir(full)
                size = stat.st_size
                mtime = datetime.datetime.fromtimestamp(stat.st_mtime).strftime(
                    "%Y-%m-%d %H:%M"
                )
                icon = "📁" if is_dir else "📄"
                items.append(f"{icon} {name}  ({mtime}, {_format_size(size)})")
            except Exception:
                items.append(f"📄 {name}")

        result = f"文件夹: {path}\n共 {len(items)} 项\n\n" + "\n".join(items)
        return result[:3000]
    except PermissionError:
        return f"错误：无权限访问「{path}」"
    except Exception as e:
        return f"错误：{str(e)}"


@register_tool(
    name="search_files",
    description="按文件名搜索电脑上的文件（支持通配符）",
    parameters={
        "pattern": {"type": "string", "description": "文件名关键词，如 *.py 或 report"},
        "folder": {
            "type": "string",
            "description": "搜索起始文件夹，默认 Desktop",
        },
        "max_results": {
            "type": "integer",
            "description": "最大返回结果数，默认 10",
        },
    },
    category="filesystem",
    timeout=30,
)
def search_files(pattern: str, folder: str = "", max_results: int = 10) -> str:
    if not folder:
        folder = os.path.expanduser("~\\Desktop")
    if not os.path.isdir(folder):
        return f"错误：文件夹不存在「{folder}」"

    try:
        results = []
        for root, dirs, files in os.walk(folder):
            # 跳过隐藏目录
            dirs[:] = [d for d in dirs if not d.startswith(".") and not d.startswith("$")]
            for f in files:
                if pattern.lower() in f.lower() or f.lower().endswith(pattern.lower().lstrip("*")):
                    full = os.path.join(root, f)
                    try:
                        size = os.path.getsize(full)
                        results.append(f"{full} ({_format_size(size)})")
                    except Exception:
                        results.append(full)
                    if len(results) >= max_results:
                        break
            if len(results) >= max_results:
                break

        if not results:
            return f"在「{folder}」下未找到包含「{pattern}」的文件"
        return f"搜索「{pattern}」结果（前 {len(results)} 个）:\n" + "\n".join(results)
    except PermissionError:
        return "搜索过程中遇到无权限的目录（已跳过）"
    except Exception as e:
        return f"搜索出错：{str(e)}"


@register_tool(
    name="read_document",
    description="读取并解析电脑上的文档文件（支持 PDF/Word/Excel/PowerPoint/代码/文本等），PDF可提取表格",
    parameters={
        "path": {"type": "string", "description": "文件完整路径"},
        "file_path": {"type": "string", "description": "文件完整路径（path 的别名）"},
        "document_path": {"type": "string", "description": "文件完整路径（path 的别名）"},
        "filepath": {"type": "string", "description": "文件完整路径（path 的别名）"},
        "extract_tables": {
            "type": "boolean",
            "description": "是否从PDF中提取表格（默认 true），仅PDF有效",
        },
    },
    category="filesystem",
    timeout=30,
)
def read_document(path: str = "", file_path: str = "", document_path: str = "", filepath: str = "", extract_tables: bool = True) -> str:
    if file_path and not path:
        path = file_path
    if document_path and not path:
        path = document_path
    if filepath and not path:
        path = filepath
    if not path:
        return "错误：请提供文件路径"
    path = os.path.abspath(path)
    if not os.path.isfile(path):
        return f"错误：文件不存在「{path}」"

    if not is_supported(path):
        return f"错误：不支持的文件类型「{os.path.splitext(path)[1]}」"

    # 先检查文件大小（超过20MB不处理）
    size = os.path.getsize(path)
    if size > 20 * 1024 * 1024:
        return f"错误：文件过大（{_format_size(size)}），暂不支持超过 20MB 的文件"

    content = parse_document(path, extract_tables=extract_tables)
    if content is None:
        return f"错误：无法解析文件「{os.path.basename(path)}」"

    ext = os.path.splitext(path)[1].lower()
    meta = get_file_metadata(path)
    summary = (
        f"文件: {os.path.basename(path)}\n"
        f"路径: {path}\n"
        f"大小: {_format_size(meta['size'])}\n"
        f"修改时间: {datetime.datetime.fromtimestamp(meta['modified']).strftime('%Y-%m-%d %H:%M')}\n"
        f"类型: {meta['category']}\n"
    )

    # PPTX显示幻灯片数
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(path)
            summary += f"幻灯片: {len(prs.slides)} 张\n"
        except Exception:
            pass
    # PDF显示页数
    elif ext == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(path)
            summary += f"页数: {len(reader.pages)}\n"
        except Exception:
            pass

    summary += "---\n"

    # 返回内容前 6000 字符（对过大文件做截断）
    max_content = 6000
    if len(content) > max_content:
        content = content[:max_content] + f"\n\n...（文件较长，仅显示前 {max_content} 字符）"

    return summary + content


# ──────────── 知识库工具 ────────────


@register_tool(
    name="query_knowledge",
    description="从知识库中检索与问题相关的文档片段（基于语义相似度）",
    parameters={
        "query": {"type": "string", "description": "搜索关键词或问题"},
        "n_results": {
            "type": "integer",
            "description": "返回结果数量，默认 5",
        },
    },
    category="knowledge",
    timeout=15,
)
def query_knowledge(query: str, n_results: int = 5) -> str:
    vs = get_vector_store()
    if vs.count() == 0:
        return "知识库为空，请先索引文件到知识库"

    try:
        results = vs.search(query, n_results=n_results)
        if not results:
            return f"知识库中未找到与「{query}」相关的内容"

        lines = [f"知识库检索「{query}」结果（共 {len(results)} 条）:\n"]
        for i, r in enumerate(results, 1):
            meta = r.get("metadata", {})
            source = meta.get("file_name", "未知来源")
            file_path = meta.get("file_path", "")
            score = r.get("score", 0)
            relevance = "高" if score < 0.3 else "中" if score < 0.6 else "低"
            lines.append(
                f"[{i}] 📄 {source}（相关度: {relevance}）\n"
                f"    路径: {file_path}\n"
                f"    内容: {r['content'][:500]}"
            )

        return "\n\n".join(lines)
    except Exception as e:
        return f"知识库检索出错：{str(e)}"


@register_tool(
    name="index_knowledge",
    description="将文件或文件夹添加到知识库用于后续检索",
    parameters={
        "path": {"type": "string", "description": "文件或文件夹路径"},
        "recursive": {
            "type": "boolean",
            "description": "是否递归扫描子文件夹，默认 true",
        },
    },
    category="knowledge",
    timeout=300,
)
def index_knowledge(path: str, recursive: bool = True) -> str:
    path = os.path.abspath(path)

    if not os.path.exists(path):
        return f"错误：路径不存在「{path}」"

    indexer = get_indexer()

    if os.path.isfile(path):
        if not is_supported(path):
            return f"错误：不支持的文件类型「{os.path.splitext(path)[1]}」"
        result = indexer.index_file(path)
        if result:
            return (
                f"✅ 已索引文件: {result['file_name']}\n"
                f"路径: {result['file_path']}\n"
                f"分块数: {result['chunk_count']}"
            )
        return f"❌ 文件索引失败: {os.path.basename(path)}"

    elif os.path.isdir(path):
        result = indexer.index_directory(path, recursive=recursive)
        return (
            f"✅ 文件夹索引完成\n"
            f"路径: {path}\n"
            f"总计: {result['total']} 个文件\n"
            f"成功: {result['success']} 个\n"
            f"失败: {result['failed']} 个\n"
            f"递归扫描: {'是' if recursive else '否'}"
        )

    return f"错误：路径无效「{path}」"


@register_tool(
    name="knowledge_stats",
    description="查看知识库统计信息（文档数量、向量数量、文件类型分布）",
    parameters={},
    category="knowledge",
    timeout=5,
)
def knowledge_stats() -> str:
    vs = get_vector_store()
    db = get_db()

    doc_stats = db.get_document_stats()
    vector_count = vs.count()

    lines = [
        "📊 知识库统计",
        f"文档总数: {doc_stats['total_documents']}",
        f"文本块总数: {doc_stats['total_chunks']}",
        f"向量总数: {vector_count}",
    ]

    if doc_stats["by_type"]:
        lines.append("\n文件类型分布:")
        for item in doc_stats["by_type"]:
            lines.append(f"  {item['file_type']}: {item['cnt']} 个")

    return "\n".join(lines)


# ──────────── 文档生成工具 ────────────


@register_tool(
    name="create_document",
    description="创建 Word 文档，可包含标题、段落、表格、页眉页脚、页码。count 参数直接指定生成份数，一次调用全部完成，无需分批",
    parameters={
        "title": {"type": "string", "description": "文档标题"},
        "content": {"type": "string", "description": "文档正文内容（支持多段，用 \\n 分隔）"},
        "table_data": {
            "type": "string",
            "description": "表格数据，JSON 格式的二维数组，如 [[\"列1\",\"列2\"],[\"a\",\"b\"]]",
        },
        "header_text": {
            "type": "string",
            "description": "页眉文字（可选）",
        },
        "footer_text": {
            "type": "string",
            "description": "页脚文字（可选）",
        },
        "show_page_numbers": {
            "type": "boolean",
            "description": "是否显示页码文字（如\"第 1 页\"），默认 false",
        },
        "save_path": {
            "type": "string",
            "description": "保存路径（含文件名），如 D:\\work\\报告.docx，默认保存到桌面",
        },
        "count": {
            "type": "integer",
            "description": "生成份数，大于1时会批量生成多份（标题后自动加编号），默认1",
        },
    },
    category="office",
    timeout=600,
)
def create_document(
    title: str = "",
    content: str = "",
    table_data: str = "",
    header_text: str = "",
    footer_text: str = "",
    show_page_numbers: bool = False,
    save_path: str = "",
    count: int = 1,
) -> str:
    try:
        from docx import Document
        from docx.shared import Inches, Pt, Cm, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return "错误：未安装 python-docx 库"

    if count < 1:
        count = 1

    need_numbering = (count > 1)

    saved_files = []
    for i in range(1, count + 1):
        doc = Document()

        # ---- 页眉 ----
        if header_text:
            header = doc.sections[0].header
            hp = header.paragraphs[0]
            hp.text = header_text
            hp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in hp.runs:
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(128, 128, 128)

        # ---- 页脚（页码仅文字，不走XML域避免文件损坏） ----
        if show_page_numbers or footer_text:
            footer = doc.sections[0].footer
            fp = footer.paragraphs[0]
            fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            parts = []
            if footer_text:
                parts.append(footer_text)
            if show_page_numbers:
                parts.append("第 {PAGE} 页")
            fp.add_run(" — ".join(parts)).font.size = Pt(9)

        # ---- 标题（批量时加编号） ----
        doc_title = title
        if need_numbering:
            doc_title = f"{title}_{i}" if title else f"文档_{i}"

        if doc_title:
            heading = doc.add_heading(doc_title, level=1)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # ---- 正文 ----
        if content:
            for para_text in content.split("\n"):
                para_text = para_text.strip()
                if not para_text:
                    continue

                if para_text.startswith("## "):
                    p = doc.add_heading(para_text[3:].strip(), level=2)
                elif para_text.startswith("### "):
                    p = doc.add_heading(para_text[4:].strip(), level=3)
                elif para_text.startswith("- ") or para_text.startswith("* "):
                    p = doc.add_paragraph(para_text[2:].strip(), style="List Bullet")
                elif para_text.startswith("1. ") or re.match(r"^\d+\.\s", para_text):
                    p = doc.add_paragraph(re.sub(r"^\d+\.\s", "", para_text), style="List Number")
                else:
                    p = doc.add_paragraph(para_text)

        # ---- 表格 ----
        if table_data:
            try:
                rows_data = json.loads(table_data)
                if rows_data and isinstance(rows_data, list):
                    table = doc.add_table(rows=len(rows_data), cols=len(rows_data[0]))
                    table.style = "Light Grid Accent 1"
                    for cj, cell_text in enumerate(rows_data[0]):
                        cell = table.cell(0, cj)
                        cell.text = str(cell_text)
                        for para in cell.paragraphs:
                            for run in para.runs:
                                run.bold = True
                    for ri in range(1, len(rows_data)):
                        for cj, cell_text in enumerate(rows_data[ri]):
                            table.cell(ri, cj).text = str(cell_text)
            except (json.JSONDecodeError, IndexError):
                pass

        # 保存
        file_path = save_path
        if not file_path:
            desktop = os.path.expanduser("~\\Desktop")
            safe = title.replace("/", "_").replace("\\", "_") if title else "未命名文档"
            if need_numbering:
                file_path = os.path.join(desktop, f"{safe}_{i}.docx")
            else:
                file_path = os.path.join(desktop, f"{safe}.docx")
        elif need_numbering:
            base, ext = os.path.splitext(file_path)
            file_path = f"{base}_{i}{ext}"

        os.makedirs(os.path.dirname(os.path.abspath(file_path)), exist_ok=True)
        doc.save(file_path)
        saved_files.append(file_path)

    if count == 1:
        return f"文档已保存至: {saved_files[0]}"
    return f"成功生成 {len(saved_files)} 份文档:\n" + "\n".join(f"  {f}" for f in saved_files)


# ──── end of create_document ────


@register_tool(
    name="create_table",
    description="创建表格数据——默认生成真正的 Excel 文件（.xlsx），支持公式、图表、多工作表，也可输出 markdown/csv/html",
    parameters={
        "headers": {"type": "string", "description": "表头，逗号分隔，如 姓名,年龄,城市"},
        "rows": {
            "type": "string",
            "description": "数据行，每行用 | 分隔，如 张三,25,北京|李四,30,上海",
        },
        "format": {
            "type": "string",
            "description": "输出格式：xlsx（默认，生成Excel文件）、markdown（聊天框内显示）、csv 或 html",
        },
        "sheet_name": {
            "type": "string",
            "description": "工作表名称（仅 xlsx 格式），默认 数据表",
        },
        "chart_type": {
            "type": "string",
            "description": "图表类型：bar（柱状图）、line（折线图）、pie（饼图），仅 xlsx 格式可选",
        },
        "chart_columns": {
            "type": "string",
            "description": "图表用的列名或列号（逗号分隔），默认使用第2列及以后的数值列",
        },
        "formulas": {
            "type": "string",
            "description": "额外公式列，JSON 格式：{\"合计\":\"SUM(B2:B{last})\",\"平均\":\"AVERAGE(B2:B{last})\"}",
        },
        "save_path": {
            "type": "string",
            "description": "Excel 保存路径，如 D:\\work\\报表.xlsx，默认保存到桌面",
        },
    },
    category="office",
    timeout=15,
)
def create_table(
    headers: str = "",
    rows: str = "",
    format: str = "xlsx",
    sheet_name: str = "数据表",
    chart_type: str = "",
    chart_columns: str = "",
    formulas: str = "",
    save_path: str = "",
) -> str:
    header_list = [h.strip() for h in headers.split(",") if h.strip()]
    if not header_list:
        return "错误：请提供表头"

    row_list = []
    for row in rows.split("|"):
        row = row.strip()
        if row:
            cells = [c.strip() for c in row.split(",")]
            if len(cells) == len(header_list):
                row_list.append(cells)
            else:
                row_list.append(cells + [""] * (len(header_list) - len(cells)))

    if format == "csv":
        output = io.StringIO()
        writer = csv.writer(output)
        writer.writerow(header_list)
        writer.writerows(row_list)
        return output.getvalue()

    elif format == "xlsx":
        if not save_path:
            desktop = os.path.expanduser("~\\Desktop")
            save_path = os.path.join(desktop, "表格数据.xlsx")
        os.makedirs(os.path.dirname(os.path.abspath(save_path)), exist_ok=True)

        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill, Border, Side, numbers
            from openpyxl.chart import BarChart, LineChart, PieChart, Reference
            from openpyxl.utils import get_column_letter

            wb = Workbook()
            ws = wb.active
            ws.title = sheet_name or "数据表"

            # 样式定义
            header_font = Font(bold=True, color="FFFFFF", size=11, name="微软雅黑")
            header_fill = PatternFill(start_color="4F6EF7", end_color="4F6EF7", fill_type="solid")
            header_align = Alignment(horizontal="center", vertical="center")
            data_font = Font(size=10, name="微软雅黑")
            thin_border = Border(
                left=Side(style="thin", color="D9D9D9"),
                right=Side(style="thin", color="D9D9D9"),
                top=Side(style="thin", color="D9D9D9"),
                bottom=Side(style="thin", color="D9D9D9"),
            )
            alt_fill = PatternFill(start_color="F5F5F5", end_color="F5F5F5", fill_type="solid")

            # 写入表头
            for j, h in enumerate(header_list, 1):
                cell = ws.cell(row=1, column=j, value=h)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_align
                cell.border = thin_border

            # 写入数据
            for i, row in enumerate(row_list, 2):
                for j, val in enumerate(row, 1):
                    cell = ws.cell(row=i, column=j, value=val)
                    cell.font = data_font
                    cell.border = thin_border
                    cell.alignment = Alignment(vertical="center")
                    # 隔行变色
                    if i % 2 == 0:
                        cell.fill = alt_fill

            # ---- 公式列 ----
            if formulas:
                try:
                    formula_config = json.loads(formulas)
                    last_row = len(row_list) + 1
                    for col_name, formula_tpl in formula_config.items():
                        j = len(header_list) + 1
                        # 表头
                        cell = ws.cell(row=1, column=j, value=col_name)
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = header_align
                        cell.border = thin_border
                        # 公式
                        formula_str = formula_tpl.replace("{last}", str(last_row))
                        for ri in range(2, last_row + 1):
                            f = formula_str.replace("{row}", str(ri))
                            cell = ws.cell(row=ri, column=j, value=f)
                            cell.font = data_font
                            cell.border = thin_border
                        header_list.append(col_name)
                except (json.JSONDecodeError, Exception) as e:
                    pass

            # ---- 图表 ----
            if chart_type and chart_type in ("bar", "line", "pie"):
                data_start_row = 2
                data_end_row = len(row_list) + 1

                # 确定图表列
                if chart_columns:
                    chart_cols = [c.strip() for c in chart_columns.split(",")]
                    col_indices = []
                    for cc in chart_cols:
                        if cc.isdigit():
                            col_indices.append(int(cc))
                        else:
                            for idx, h in enumerate(header_list):
                                if cc.lower() in h.lower():
                                    col_indices.append(idx + 1)
                else:
                    # 默认用第2列起的数值列
                    col_indices = list(range(2, len(header_list) + 1))

                if col_indices:
                    cats = Reference(ws, min_col=1, min_row=data_start_row,
                                     max_row=data_end_row)
                    chart = None
                    if chart_type == "bar":
                        chart = BarChart()
                        chart.type = "col"
                    elif chart_type == "line":
                        chart = LineChart()
                    elif chart_type == "pie":
                        chart = PieChart()

                    if chart:
                        chart.title = f"{' '.join(header_list[:2])} 图表"
                        chart.style = 10
                        for ci in col_indices:
                            if ci <= len(header_list):
                                data = Reference(ws, min_col=ci, min_row=1,
                                                 max_row=data_end_row)
                                chart.add_data(data, titles_from_data=True)
                        if cats:
                            chart.set_categories(cats)
                        chart.shape = 4
                        ws.add_chart(chart, f"E{data_end_row + 3}")

            # 自动列宽
            for col in ws.columns:
                max_len = 0
                for cell in col:
                    if cell.value:
                        width = sum(2 if ord(c) > 127 else 1 for c in str(cell.value))
                        max_len = max(max_len, width)
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

            wb.save(save_path)
            chart_msg = f" + {chart_type}图" if chart_type else ""
            formula_msg = " + 公式" if formulas else ""
            return (
                f"Excel 文件已生成并保存！\n"
                f"文件路径: {save_path}\n"
                f"包含 {len(header_list)} 列, {len(row_list)} 行数据"
                f"{chart_msg}{formula_msg}\n"
                f"直接用 Excel 打开即可。"
            )
        except ImportError as e:
            return f"错误：缺少库 ({e})，无法生成 Excel 文件"
        except Exception as e:
            return f"生成 Excel 失败: {str(e)}"

    elif format == "html":
        html = ['<table border="1" cellpadding="4" cellspacing="0" style="border-collapse:collapse;width:100%">']
        html.append("<thead><tr>" + "".join(f"<th>{h}</th>" for h in header_list) + "</tr></thead>")
        html.append("<tbody>")
        for row in row_list:
            html.append("<tr>" + "".join(f"<td>{c}</td>" for c in row) + "</tr>")
        html.append("</tbody></table>")
        return "\n".join(html)
    else:
        lines = ["| " + " | ".join(header_list) + " |"]
        lines.append("| " + " | ".join(["---"] * len(header_list)) + " |")
        for row in row_list:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines)


# ──────────── 数据处理工具 ────────────


@register_tool(
    name="analyze_data",
    description="分析 CSV、Excel 或表格数据，计算统计指标（总和、平均值、最大最小等）",
    parameters={
        "data_source": {
            "type": "string",
            "description": "数据来源：文件路径（支持.csv/.xlsx）或直接输入表格内容",
        },
        "column": {
            "type": "string",
            "description": "要分析的列名或列号（从0开始），默认分析所有数值列",
        },
        "operations": {
            "type": "string",
            "description": "分析操作，逗号分隔：sum,avg,max,min,count,std，默认全部",
        },
    },
    category="data",
    timeout=30,
)
def analyze_data(data_source: str, column: str = "", operations: str = "") -> str:
    rows = []

    # 尝试从文件读取
    if os.path.isfile(data_source):
        ext = os.path.splitext(data_source)[1].lower()
        try:
            if ext == ".csv":
                with open(data_source, "r", encoding="utf-8") as f:
                    reader = csv.reader(f)
                    rows = list(reader)
            elif ext == ".xlsx":
                from openpyxl import load_workbook

                wb = load_workbook(data_source, read_only=True, data_only=True)
                ws = wb.active
                for row in ws.iter_rows(values_only=True):
                    rows.append([str(c) if c is not None else "" for c in row])
            else:
                return f"不支持的文件格式: {ext}"
        except Exception as e:
            return f"读取文件失败: {e}"
    else:
        # 尝试将输入解析为表格（每行一行，逗号或制表符分隔）
        for line in data_source.strip().split("\n"):
            line = line.strip()
            if line:
                rows.append([c.strip() for c in line.replace("\t", ",").split(",")])

    if not rows or len(rows) < 2:
        return "数据不足，至少需要表头和一行数据"

    headers = rows[0]
    data_rows = rows[1:]

    # 确定要分析的列
    col_indices = []
    if column:
        cols = [c.strip() for c in column.split(",")]
        for col in cols:
            if col.isdigit():
                col_indices.append(int(col))
            else:
                for i, h in enumerate(headers):
                    if col.lower() in h.lower():
                        col_indices.append(i)
        col_indices = list(set(col_indices))
    else:
        # 自动检测数值列
        for i in range(len(headers)):
            vals = []
            for r in data_rows:
                if i < len(r):
                    try:
                        float(r[i].replace(",", "").replace("¥", "").replace("$", ""))
                        vals.append(True)
                    except ValueError:
                        vals.append(False)
            if any(vals) and len(vals) > len(data_rows) * 0.5:
                col_indices.append(i)

    if not col_indices:
        return "未找到可分析的数值列"

    # 解析操作
    op_list = [o.strip().lower() for o in operations.split(",") if o.strip()]
    if not op_list:
        op_list = ["sum", "avg", "max", "min", "count"]

    op_names = {
        "sum": "总和",
        "avg": "平均值",
        "max": "最大值",
        "min": "最小值",
        "count": "计数",
        "std": "标准差",
    }

    results = [f"📊 数据分析结果\n数据行数: {len(data_rows)}\n"]
    for ci in col_indices:
        col_name = headers[ci] if ci < len(headers) else f"列{ci}"
        values = []
        for r in data_rows:
            if ci < len(r):
                try:
                    v = float(r[ci].replace(",", "").replace("¥", "").replace("$", ""))
                    values.append(v)
                except ValueError:
                    pass

        if not values:
            continue

        results.append(f"\n▶ {col_name} ({len(values)} 个数值):")
        for op in op_list:
            op = op.lower()
            if op == "sum":
                results.append(f"  总和: {sum(values):.2f}")
            elif op == "avg":
                results.append(f"  平均值: {statistics.mean(values):.2f}")
            elif op == "max":
                results.append(f"  最大值: {max(values):.2f}")
            elif op == "min":
                results.append(f"  最小值: {min(values):.2f}")
            elif op == "count":
                results.append(f"  计数: {len(values)}")
            elif op == "std":
                results.append(f"  标准差: {statistics.stdev(values):.2f}" if len(values) > 1 else "  标准差: 数据不足")

    return "\n".join(results)


@register_tool(
    name="format_data",
    description="格式化数据，支持排序、筛选、转置等操作",
    parameters={
        "data": {"type": "string", "description": "输入数据，CSV 格式或文件路径"},
        "operation": {
            "type": "string",
            "description": "操作类型：sort（排序）、filter（筛选）、transpose（转置）、head（预览前N行）",
        },
        "params": {
            "type": "string",
            "description": "操作参数，如 sort:列号/asc 或 filter:列号>100 或 head:5",
        },
    },
    category="data",
    timeout=15,
)
def format_data(data: str, operation: str = "head", params: str = "5") -> str:
    rows = []
    if os.path.isfile(data):
        try:
            with open(data, "r", encoding="utf-8") as f:
                reader = csv.reader(f)
                rows = list(reader)
        except Exception:
            return f"无法读取文件: {data}"
    else:
        for line in data.strip().split("\n"):
            line = line.strip()
            if line:
                rows.append([c.strip() for c in line.split(",")])

    if not rows:
        return "无数据"

    op = operation.strip().lower()
    result_rows = list(rows)  # copy

    try:
        if op == "head":
            n = int(params) if params.isdigit() else 5
            result_rows = rows[: min(n + 1, len(rows))]

        elif op == "sort":
            parts = params.split("/")
            col = int(parts[0]) if parts[0].isdigit() else 0
            reverse = len(parts) > 1 and parts[1].strip().lower() == "desc"
            header = result_rows[0]
            data_rows = result_rows[1:]
            data_rows.sort(
                key=lambda r: (
                    float(r[col].replace(",", ""))
                    if col < len(r) and r[col].replace(",", "").replace(".", "").isdigit()
                    else r[col] if col < len(r) else ""
                ),
                reverse=reverse,
            )
            result_rows = [header] + data_rows

        elif op == "filter":
            parts = params.split(">") if ">" in params else params.split("<") if "<" in params else [params, ""]
            if len(parts) != 2:
                return "过滤格式: 列号>值 或 列号=值"
            col_str, val_str = parts
            col = int(col_str.strip()) if col_str.strip().isdigit() else 0
            is_gt = ">" in params
            header = result_rows[0]
            data_rows = result_rows[1:]
            filtered = []
            for r in data_rows:
                if col < len(r):
                    cell = r[col].replace(",", "").strip()
                    try:
                        cell_val = float(cell)
                        threshold = float(val_str.strip())
                        if (is_gt and cell_val > threshold) or (not is_gt and cell_val < threshold):
                            filtered.append(r)
                    except ValueError:
                        if (is_gt and cell > val_str.strip()) or (not is_gt and cell < val_str.strip()):
                            filtered.append(r)
            result_rows = [header] + filtered if filtered else [header] + ["(无匹配行)"]

        elif op == "transpose":
            result_rows = list(map(list, zip(*rows)))

    except Exception as e:
        return f"操作失败: {str(e)}"

    # 格式化为表格
    out = io.StringIO()
    writer = csv.writer(out)
    writer.writerows(result_rows)
    return out.getvalue()


# ──────────── 文件/文件夹创建工具 ────────────


@register_tool(
    name="create_folder",
    description="创建文件夹（如果父路径不存在会自动创建）",
    parameters={
        "path": {"type": "string", "description": "要创建的文件夹路径，如 D:\\work\\报告 或 C:\\Users\\xxx\\Desktop\\新文件夹"},
    },
    category="filesystem",
    timeout=10,
)
def create_folder(path: str) -> str:
    path = os.path.abspath(path)
    try:
        os.makedirs(path, exist_ok=True)
        return f"✅ 文件夹已创建: {path}"
    except PermissionError:
        return f"错误：无权限创建文件夹「{path}」"
    except Exception as e:
        return f"创建文件夹失败: {str(e)}"


# ──────────── 文件删除工具 ────────────


@register_tool(
    name="delete_file",
    description="删除电脑上的文件（不可恢复！请先确认用户同意再执行）",
    parameters={
        "path": {"type": "string", "description": "要删除的文件完整路径"},
        "confirm": {
            "type": "boolean",
            "description": "是否确认删除，必须为 true 才能执行",
        },
    },
    category="filesystem",
    timeout=10,
)
def delete_file(path: str, confirm: bool = False) -> str:
    import shutil
    import stat

    path = os.path.abspath(path)

    if not confirm:
        return f"⚠️ 请确认是否要删除「{path}」？如需删除，请设置 confirm=true"

    if not os.path.exists(path):
        return f"错误：文件不存在「{path}」"
    if os.path.isdir(path):
        return f"错误：不支持删除文件夹，请手动删除文件夹「{path}」"

    # 检查是否只读，自动去除只读属性
    try:
        mode = os.stat(path).st_mode
        if not mode & stat.S_IWRITE:
            os.chmod(path, mode | stat.S_IWRITE)
    except Exception:
        pass

    try:
        os.remove(path)
        return f"✅ 文件已删除: {path}"
    except PermissionError:
        # 判断是否被其他进程占用
        try:
            with open(path, 'a'):
                pass
        except PermissionError:
            return (f"错误：无权限删除文件「{path}」\n"
                    f"可能原因：文件被其他程序打开（如Excel/WPS/文本编辑器），请关闭后再试。")
        except OSError:
            return f"错误：文件「{path}」被其他进程占用，请关闭相关程序后重试。"
        return f"错误：无权限删除文件「{path}」，请检查文件权限设置。"
    except Exception as e:
        return f"删除失败: {str(e)}"


@register_tool(
    name="remove_from_knowledge",
    description="从知识库中删除已索引的文档（按文件名或路径匹配）",
    parameters={
        "name": {
            "type": "string",
            "description": "要删除的文件名或路径关键词，如 report.docx 或 D:\\work\\data.xlsx",
        },
    },
    category="knowledge",
    timeout=10,
)
def remove_from_knowledge(name: str) -> str:
    from core.database.db import get_db
    from core.rag.vector_store import get_vector_store

    db = get_db()
    vs = get_vector_store()
    docs = db.list_documents()

    matched = [d for d in docs if name.lower() in d["file_name"].lower() or name.lower() in d["file_path"].lower()]

    if not matched:
        return f"知识库中未找到包含「{name}」的文档"

    deleted = []
    for doc in matched:
        try:
            vs.delete_document(doc["id"])
            db.remove_document(doc["id"])
            deleted.append(doc["file_name"])
        except Exception as e:
            pass

    if deleted:
        return f"已从知识库删除 {len(deleted)} 个文档:\n" + "\n".join(f"- {f}" for f in deleted)
    return "删除操作未生效"


# ──────────── 输出验证工具 ────────────


def verify_output(path: str = "", expect_rows: int = 0, check_columns: str = "",
                  sheet_name: str = "", check_sum: str = "") -> str:
    """验证输出文件（Excel/CSV）的数据内容完整性。

    Args:
        path: 文件路径
        expect_rows: 期望的数据行数（0=不检查，排除表头行）
        check_columns: 需要存在的列名，逗号分隔
        sheet_name: 要检查的 sheet 名（Excel 专用，空=全部 sheet）
        check_sum: 需要验证合计的格式 "列名=期望值"，如 "金额=1000" 或 "金额=-"（只检查能求和）

    Returns:
        PASS/FAIL 报告
    """
    import os
    if not path or not os.path.exists(path):
        return f"FAIL: 文件不存在: {path}"

    ext = os.path.splitext(path)[1].lower()
    reports = []

    try:
        if ext in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)

            # 确定要检查的 sheet
            if sheet_name:
                sheets = [s for s in wb.sheetnames if s == sheet_name]
                if not sheets:
                    wb.close()
                    return f"FAIL: 未找到 Sheet '{sheet_name}'，可用: {wb.sheetnames}"
            else:
                sheets = wb.sheetnames

            for sn in sheets:
                ws = wb[sn]
                sr = _verify_sheet(ws, sn, expect_rows, check_columns, check_sum)
                reports.append(sr)

            wb.close()

        elif ext == ".csv":
            import csv
            with open(path, "r", encoding="utf-8-sig") as f:
                reader = csv.reader(f)
                headers = next(reader)
                columns = [h.strip() for h in headers]
                rows = []
                for row in reader:
                    rows.append(row)

            data_rows = len(rows)
            nulls = []
            issues = []

            for ci, col_name in enumerate(columns):
                null_count = sum(1 for r in rows if ci < len(r) and not r[ci].strip())
                if null_count > 0:
                    nulls.append(f"{col_name}({null_count}个)")

            if expect_rows > 0 and data_rows != expect_rows:
                issues.append(f"行数不符: 期望 {expect_rows}，实际 {data_rows}")
            if check_columns:
                req_cols = [c.strip() for c in check_columns.split(",")]
                for rc in req_cols:
                    if rc not in columns:
                        issues.append(f"缺少列: {rc}")
            if nulls:
                issues.append(f"空值: {', '.join(nulls)}")

            summary = f"  Sheet: {os.path.basename(path)} | 行数: {data_rows} | 列数: {len(columns)}"
            if columns:
                summary += f" | 列: {', '.join(columns[:8])}"

            if issues:
                reports.append(f"FAIL [{os.path.basename(path)}]: {'; '.join(issues)}\n{summary}")
            else:
                reports.append(f"PASS [{os.path.basename(path)}]\n{summary}")

        else:
            return f"SKIP: 不支持的文件类型 {ext}，仅支持 .xlsx/.xls/.csv"

    except Exception as e:
        return f"FAIL: 文件读取失败 - {e}"

    failed = [r for r in reports if r.startswith("FAIL")]
    if failed:
        return "\n\n".join(failed)
    return "PASS\n" + "\n\n".join(reports)


def _verify_sheet(ws, sheet_name: str, expect_rows: int,
                  check_columns: str, check_sum: str) -> str:
    """验证单个 Excel Sheet 的数据完整性。"""
    rows_iter = list(ws.iter_rows(values_only=True))
    if not rows_iter:
        return f"SKIP [{sheet_name}]: 空 Sheet"

    # 找第一个非空行作为表头行
    header_row_idx = None
    for i, row in enumerate(rows_iter):
        vals = [v for v in row if v is not None]
        if len(vals) >= 2:  # 至少有2个非空列才认为是表头
            header_row_idx = i
            break

    if header_row_idx is None:
        return f"SKIP [{sheet_name}]: 未找到表头行"

    columns = [str(c or "").strip() for c in rows_iter[header_row_idx]]
    columns = [c for c in columns if c]  # 去掉空列名

    # 数据行（表头之后的非空行）
    data_rows = []
    for row in rows_iter[header_row_idx + 1:]:
        vals = [v for v in row if v is not None and str(v).strip()]
        if vals:  # 至少有一个非空值才认为是数据行
            data_rows.append(row)

    data_count = len(data_rows)
    issues = []

    # 行数检查
    if expect_rows > 0 and data_count != expect_rows:
        issues.append(f"行数不符: 期望 {expect_rows}，实际 {data_count}")

    # 列名检查
    if check_columns:
        req_cols = [c.strip() for c in check_columns.split(",")]
        for rc in req_cols:
            if rc not in columns:
                issues.append(f"缺少列: {rc}")

    # 空值检查（只检查表头定义的列）
    col_null_map = {}
    for ri, row in enumerate(data_rows):
        for ci in range(min(len(row), len(columns))):
            col_name = columns[ci]
            if not col_name:
                continue
            val = row[ci]
            if val is None or (isinstance(val, str) and not val.strip()):
                col_null_map[col_name] = col_null_map.get(col_name, 0) + 1
    null_cols = [f"{c}({n}个)" for c, n in col_null_map.items() if n > 0]
    if null_cols:
        issues.append(f"空值: {', '.join(null_cols)}")

    # 合计检查
    if check_sum:
        for part in check_sum.split(";"):
            part = part.strip()
            if "=" not in part:
                # 纯列名 → 自动求和
                col_name = part
                try:
                    ci = columns.index(col_name)
                    total = 0
                    for row in data_rows:
                        if ci < len(row):
                            v = row[ci]
                            if v is not None:
                                total += float(v)
                    issues.append(f"合计[{col_name}]: {total:.2f}")
                except (ValueError, IndexError):
                    issues.append(f"合计检查: 列 '{col_name}' 不存在")
            else:
                col_name, expected = part.split("=", 1)
                col_name = col_name.strip()
                expected = expected.strip()
                try:
                    ci = columns.index(col_name)
                    total = 0
                    for row in data_rows:
                        if ci < len(row):
                            v = row[ci]
                            if v is not None:
                                total += float(v)
                    # 期望值是 "-" 或空 → 自动求和（不校验）
                    if expected in ("", "-"):
                        issues.append(f"自动合计[{col_name}]: {total:.2f}")
                    else:
                        exp_val = float(expected)
                        if abs(total - exp_val) > 0.01:
                            issues.append(f"合计[{col_name}]不符: 期望 {exp_val:.2f}，实际 {total:.2f}")
                        else:
                            issues.append(f"合计[{col_name}]验证通过: {total:.2f}")
                except ValueError:
                    issues.append(f"合计[{col_name}]求和失败: 非数值列?")
                except IndexError:
                    issues.append(f"合计检查: 列 '{col_name}' 不存在")

    summary = f"  Sheet: {sheet_name} | 数据行: {data_count} | 列数: {len(columns)}"
    if columns:
        summary += f" | 列: {', '.join(columns[:8])}"

    if issues:
        return f"FAIL [{sheet_name}]: {'; '.join(issues)}\n{summary}"
    return f"PASS [{sheet_name}]\n{summary}"


@register_tool(
    name="create_presentation",
    description="创建 PowerPoint 演示文稿（.pptx），支持图表/表格/图片/多版式/主题/模板/备注/切换效果",
    parameters={
        "title": {"type": "string", "description": "演示文稿标题（封面页用）"},
        "slides": {
            "type": "string",
            "description": "幻灯片 JSON 数组，支持:\n"
                           "- 基础: {\"title\":\"页标题\",\"content\":\"正文(\\n换行, ##=小标题)\",\"bullet\":true}\n"
                           "- 版式: layout=\"content\"|\"bullet\"|\"two_column\"|\"image_right\"|\"image_left\"|\"section\"|\"chart\"|\"table\"\n"
                           "- 图表: chart={\"type\":\"column\"|\"bar\"|\"line\"|\"pie\"|\"area\", \"categories\":[], \"series\":[]}\n"
                           "- 表格: table={\"headers\":[], \"rows\":[]}\n"
                           "- 图片: image=\"C:/photo.png\"\n"
                           "- 备注: notes=\"演讲者备注\"\n"
                           "- 切换: transition=\"fade\"|\"push\"|\"wipe\"|\"zoom\"",
        },
        "save_path": {
            "type": "string",
            "description": "保存路径（含文件名），如 D:\\work\\演示.pptx，默认保存到桌面",
        },
        "template_path": {
            "type": "string",
            "description": "模板文件路径 (.potx/.pptx)，不指定则用默认空白模板",
        },
        "theme": {
            "type": "string",
            "description": "主题配色: default(科吉蓝) | modern(红) | minimal(青) | dark(暗色) | nature(绿) | warm(暖橙)",
        },
    },
    category="office",
    timeout=60,
)
def create_presentation(title: str = "", slides: str = "", save_path: str = "",
                        template_path: str = "", theme: str = "default") -> str:
    from core.presentation import create_presentation as _build
    return _build(title=title, slides=slides, save_path=save_path,
                  template_path=template_path, theme=theme)


# ──────────── 工具函数 ────────────


def _format_size(size: int) -> str:
    for unit in ["B", "KB", "MB", "GB"]:
        if size < 1024:
            return f"{size:.1f}{unit}"
        size /= 1024
    return f"{size:.1f}TB"
