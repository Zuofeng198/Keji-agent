"""企业级 PPT 生成引擎（视觉增强版）

设计理念：
  - 内容以「卡片」形式呈现，白色卡片 + 阴影，视觉层次分明
  - 统一的中文字体排版（Microsoft YaHei），专业间距体系
  - 精致装饰元素（左侧色条、装饰线、几何元素）
  - 所有主题统一遵守：浅色背景 + 白色卡片 + 主题色点缀

支持：
  - 6 套内置主题配色
  - 多种版式：内容/两栏/图文/章节页/要点
  - 图表：柱状/折线/饼图/面积
  - 表格（主题色表头 + 交替行）
  - 图片插入（自动缩放）
  - 演讲者备注
  - 切换动画
  - 模板加载（.potx/.pptx）
"""

import json
import os
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from lxml import etree


class PresentationBuilder:
    """演示文稿构建器（视觉增强版）"""

    THEMES = {
        "default": {
            "name": "科吉默认",
            "primary": RGBColor(0x1E, 0x27, 0x61),
            "accent": RGBColor(0x4F, 0x6E, 0xF7),
            "white": RGBColor(0xFF, 0xFF, 0xFF),
            "dark": RGBColor(0x2C, 0x3E, 0x50),
            "gray": RGBColor(0x7F, 0x8C, 0x9B),
            "bg_light": RGBColor(0xF5, 0xF6, 0xFA),
            "bg_dark": RGBColor(0x1E, 0x27, 0x61),
            "bg_card": RGBColor(0xFF, 0xFF, 0xFF),
            "card_border": RGBColor(0xE8, 0xEC, 0xF1),
            "chart_colors": [
                RGBColor(0x4F, 0x6E, 0xF7),
                RGBColor(0xE7, 0x4C, 0x3C),
                RGBColor(0x27, 0xAE, 0x60),
                RGBColor(0xF3, 0x9C, 0x12),
                RGBColor(0x9B, 0x59, 0xB6),
                RGBColor(0x1A, 0xBC, 0x9C),
            ],
        },
        "modern": {
            "name": "现代简约",
            "primary": RGBColor(0x2C, 0x3E, 0x50),
            "accent": RGBColor(0xE7, 0x4C, 0x3C),
            "white": RGBColor(0xFF, 0xFF, 0xFF),
            "dark": RGBColor(0x2C, 0x3E, 0x50),
            "gray": RGBColor(0x95, 0xA5, 0xA6),
            "bg_light": RGBColor(0xEC, 0xF0, 0xF1),
            "bg_dark": RGBColor(0x2C, 0x3E, 0x50),
            "bg_card": RGBColor(0xFF, 0xFF, 0xFF),
            "card_border": RGBColor(0xE8, 0xEC, 0xF1),
            "chart_colors": [
                RGBColor(0xE7, 0x4C, 0x3C),
                RGBColor(0x34, 0x98, 0xDB),
                RGBColor(0x2E, 0xCC, 0x71),
                RGBColor(0xF3, 0x9C, 0x12),
                RGBColor(0x9B, 0x59, 0xB6),
                RGBColor(0x1A, 0xBC, 0x9C),
            ],
        },
        "minimal": {
            "name": "极简白",
            "primary": RGBColor(0x33, 0x33, 0x33),
            "accent": RGBColor(0x00, 0x96, 0x88),
            "white": RGBColor(0xFF, 0xFF, 0xFF),
            "dark": RGBColor(0x33, 0x33, 0x33),
            "gray": RGBColor(0x99, 0x99, 0x99),
            "bg_light": RGBColor(0xFA, 0xFA, 0xFA),
            "bg_dark": RGBColor(0x33, 0x33, 0x33),
            "bg_card": RGBColor(0xFF, 0xFF, 0xFF),
            "card_border": RGBColor(0xE8, 0xE8, 0xE8),
            "chart_colors": [
                RGBColor(0x00, 0x96, 0x88),
                RGBColor(0xFF, 0x57, 0x22),
                RGBColor(0x03, 0xA9, 0xF4),
                RGBColor(0x8B, 0xC3, 0x4A),
                RGBColor(0xFF, 0x98, 0x00),
                RGBColor(0xE9, 0x1E, 0x63),
            ],
        },
        "dark": {
            "name": "暗色商务",
            "primary": RGBColor(0x1A, 0x1A, 0x2E),
            "accent": RGBColor(0x00, 0xD2, 0xFF),
            "white": RGBColor(0xFF, 0xFF, 0xFF),
            "dark": RGBColor(0xE0, 0xE0, 0xE0),
            "gray": RGBColor(0x88, 0x88, 0x88),
            "bg_light": RGBColor(0x2D, 0x2D, 0x44),
            "bg_dark": RGBColor(0x1A, 0x1A, 0x2E),
            "bg_card": RGBColor(0x2D, 0x2D, 0x44),
            "card_border": RGBColor(0x3D, 0x3D, 0x5C),
            "chart_colors": [
                RGBColor(0x00, 0xD2, 0xFF),
                RGBColor(0xFF, 0x6B, 0x6B),
                RGBColor(0xFF, 0xD9, 0x3D),
                RGBColor(0x6B, 0xCB, 0x77),
                RGBColor(0xC5, 0x6C, 0xF0),
                RGBColor(0xFF, 0x8E, 0x53),
            ],
        },
        "nature": {
            "name": "自然清新",
            "primary": RGBColor(0x27, 0xAE, 0x60),
            "accent": RGBColor(0x2E, 0xCC, 0x71),
            "white": RGBColor(0xFF, 0xFF, 0xFF),
            "dark": RGBColor(0x2C, 0x3E, 0x50),
            "gray": RGBColor(0x7F, 0x8C, 0x9D),
            "bg_light": RGBColor(0xE8, 0xF8, 0xF5),
            "bg_dark": RGBColor(0x27, 0xAE, 0x60),
            "bg_card": RGBColor(0xFF, 0xFF, 0xFF),
            "card_border": RGBColor(0xD5, 0xF0, 0xE5),
            "chart_colors": [
                RGBColor(0x27, 0xAE, 0x60),
                RGBColor(0x2E, 0xCC, 0x71),
                RGBColor(0x1A, 0xBC, 0x9C),
                RGBColor(0x16, 0xA0, 0x85),
                RGBColor(0x34, 0x98, 0xDB),
                RGBColor(0xE6, 0x7E, 0x22),
            ],
        },
        "warm": {
            "name": "暖色调",
            "primary": RGBColor(0x8B, 0x45, 0x13),
            "accent": RGBColor(0xE6, 0x7E, 0x22),
            "white": RGBColor(0xFF, 0xFF, 0xFF),
            "dark": RGBColor(0x5D, 0x40, 0x37),
            "gray": RGBColor(0xAE, 0x99, 0x8C),
            "bg_light": RGBColor(0xFD, 0xF2, 0xE9),
            "bg_dark": RGBColor(0x8B, 0x45, 0x13),
            "bg_card": RGBColor(0xFF, 0xFF, 0xFF),
            "card_border": RGBColor(0xF5, 0xE6, 0xD5),
            "chart_colors": [
                RGBColor(0xE6, 0x7E, 0x22),
                RGBColor(0xD3, 0x54, 0x00),
                RGBColor(0xF1, 0xC4, 0x0F),
                RGBColor(0x16, 0xA0, 0x85),
                RGBColor(0x8E, 0x44, 0xAD),
                RGBColor(0xC0, 0x39, 0x2B),
            ],
        },
    }

    # 切换效果 OXML 标签映射
    TRANSITION_TAGS = {
        "fade": "p:fade",
        "push": "p:push",
        "wipe": "p:wipe",
        "split": "p:split",
        "uncover": "p:uncover",
        "cover": "p:cover",
        "zoom": "p:zoom",
    }

    CHART_TYPES = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "area": XL_CHART_TYPE.AREA,
    }

    # ──────── 初始化 ────────

    def __init__(self, template_path: str = "", theme: str = "default"):
        if template_path and os.path.isfile(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)

        self.theme = self.THEMES.get(theme, self.THEMES["default"])
        self._slide_index = 0

    # ──────── 公开 API ────────

    def add_cover(self, title: str, subtitle: str = ""):
        """添加封面页（深色背景 + 装饰几何元素）"""
        t = self.theme
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, t["bg_dark"])

        sw = self.prs.slide_width
        sh = self.prs.slide_height

        # 顶部装饰条
        self._add_rect(slide, Inches(0), Inches(0), sw, Inches(0.06), t["accent"])
        # 底部装饰条
        self._add_rect(slide, Inches(0), sh - Inches(0.06), sw, Inches(0.06), t["accent"])

        # 左侧竖线装饰
        self._add_rect(slide, Inches(1.2), Inches(2.5), Inches(0.04), Inches(2.5), t["accent"])

        # 标题
        tb = slide.shapes.add_textbox(Inches(1.8), Inches(2.2),
                                      Inches(10), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        self._set_paragraph_font(p, Pt(42), t["white"], bold=True)
        p.alignment = PP_ALIGN.LEFT

        # 副标题
        if subtitle:
            tb = slide.shapes.add_textbox(Inches(1.8), Inches(4.2),
                                          Inches(10), Inches(0.8))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            self._set_paragraph_font(p, Pt(18), t["gray"])
            p.alignment = PP_ALIGN.LEFT

        # 右下角装饰方块
        self._add_rect(slide, sw - Inches(2.0), sh - Inches(2.0),
                       Inches(1.0), Inches(1.0), t["accent"])
        # 透明方块衬底
        self._add_rect(slide, sw - Inches(1.4), sh - Inches(1.4),
                       Inches(0.6), Inches(0.6),
                       RGBColor(0xFF, 0xFF, 0xFF))

        self._slide_index += 1

    def add_slide(self, data: dict):
        """根据 layout 分发到对应构建方法"""
        layout = data.get("layout", "content")
        data.setdefault("transition", None)
        data.setdefault("notes", "")

        if data.get("bullet") and layout == "content":
            layout = "bullet"

        if layout == "section":
            self._add_section(data)
        elif layout in ("image_right", "image_left"):
            self._add_image(data, layout == "image_left")
        elif layout == "two_column":
            self._add_two_column(data)
        elif layout in ("chart", "chart_only"):
            self._add_chart(data)
        elif layout in ("table", "table_only"):
            self._add_table(data)
        elif layout == "bullet":
            self._add_bullet(data)
        else:
            self._add_content(data)

        self._slide_index += 1

    def save(self, path: str):
        """保存为 pptx 文件"""
        os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
        self.prs.save(path)

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides)

    # ──────── 页面构建器 ────────

    def _add_content(self, data: dict):
        """内容页：左侧装饰色条 + 标题 + 卡片内正文"""
        t = self.theme
        slide = self._create_slide(t["bg_light"])

        title = data.get("title", f"第{self._slide_index + 1}页")
        content = data.get("content", "")

        # 左侧装饰色条
        self._add_rect(slide, Inches(0), Inches(0),
                       Inches(0.06), self.prs.slide_height, t["accent"])

        # 标题
        self._add_slide_title(slide, title)

        # 白色卡片
        card = self._add_card(slide, Inches(0.5), Inches(1.5),
                              Inches(12.333), Inches(5.6))

        # 卡片内正文
        tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.8),
                                      Inches(11.333), Inches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.1)

        lines = content.split("\n") if content else []
        for li, line in enumerate(lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            line = line.strip()
            if not line:
                continue
            if line.startswith("## "):
                p.text = line[3:]
                self._set_paragraph_font(p, Pt(20), t["accent"], bold=True)
                p.space_before = Pt(14)
                p.space_after = Pt(6)
            else:
                p.text = line
                self._set_paragraph_font(p, Pt(15), t["dark"])
                p.space_after = Pt(5)
                p.line_spacing = Pt(24)

        self._add_notes(slide, data.get("notes", ""))
        self._apply_transition(slide, data.get("transition", ""))

    def _add_bullet(self, data: dict):
        """要点页：卡片内项目符号列表"""
        t = self.theme
        slide = self._create_slide(t["bg_light"])

        title = data.get("title", f"第{self._slide_index + 1}页")
        content = data.get("content", "")

        # 左侧装饰色条
        self._add_rect(slide, Inches(0), Inches(0),
                       Inches(0.06), self.prs.slide_height, t["accent"])

        self._add_slide_title(slide, title)

        card = self._add_card(slide, Inches(0.5), Inches(1.5),
                              Inches(12.333), Inches(5.6))

        tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.8),
                                      Inches(11.333), Inches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)

        lines = content.split("\n") if content else []
        for li, line in enumerate(lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            line = line.strip()
            if not line:
                continue
            p.text = line
            self._set_paragraph_font(p, Pt(17), t["dark"])
            p.space_after = Pt(8)
            p.level = 0

        self._add_notes(slide, data.get("notes", ""))
        self._apply_transition(slide, data.get("transition", ""))

    def _add_two_column(self, data: dict):
        """两栏布局：两张并排卡片"""
        t = self.theme
        slide = self._create_slide(t["bg_light"])

        title = data.get("title", "")
        content = data.get("content", "")
        cols = content.split("\n---\n") if content else ["", ""]
        left_content = cols[0] if len(cols) > 0 else ""
        right_content = cols[1] if len(cols) > 1 else ""

        self._add_slide_title(slide, title)

        # 左卡片
        self._add_card(slide, Inches(0.5), Inches(1.5),
                       Inches(5.9), Inches(5.6))
        self._add_column_text(slide, Inches(0.8), Inches(1.7),
                              Inches(5.3), left_content, t)

        # 右卡片
        self._add_card(slide, Inches(6.9), Inches(1.5),
                       Inches(5.9), Inches(5.6))
        self._add_column_text(slide, Inches(7.2), Inches(1.7),
                              Inches(5.3), right_content, t)

        self._add_notes(slide, data.get("notes", ""))
        self._apply_transition(slide, data.get("transition", ""))

    def _add_image(self, data: dict, image_on_left: bool = False):
        """图文布局：卡片内文字 + 图片"""
        t = self.theme
        slide = self._create_slide(t["bg_light"])

        title = data.get("title", "")
        content = data.get("content", "")
        image_path = data.get("image", "")

        self._add_slide_title(slide, title)

        card = self._add_card(slide, Inches(0.5), Inches(1.5),
                              Inches(12.333), Inches(5.6))

        has_image = image_path and os.path.isfile(image_path)

        if has_image:
            text_left = Inches(1.0) if not image_on_left else Inches(7.2)
            text_width = Inches(5.5)
            img_left = Inches(7.2) if not image_on_left else Inches(1.0)

            try:
                slide.shapes.add_picture(image_path, img_left, Inches(2.0),
                                         width=Inches(5.0))
            except Exception:
                # 图片加载失败，全宽文字
                text_left = Inches(1.0)
                text_width = Inches(11.333)
        else:
            text_left = Inches(1.0)
            text_width = Inches(11.333)

        # 文字
        tb = slide.shapes.add_textbox(text_left, Inches(1.8),
                                      text_width, Inches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_right = Inches(0.2)

        lines = content.split("\n") if content else []
        for li, line in enumerate(lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            line = line.strip()
            if not line:
                continue
            p.text = line
            self._set_paragraph_font(p, Pt(16), t["dark"])
            p.space_after = Pt(5)
            p.line_spacing = Pt(26)

        self._add_notes(slide, data.get("notes", ""))
        self._apply_transition(slide, data.get("transition", ""))

    def _add_chart(self, data: dict):
        """图表页：卡片内嵌图表"""
        t = self.theme
        slide = self._create_slide(t["bg_light"])

        title = data.get("title", "")
        cd = data.get("chart", {})

        if title:
            self._add_slide_title(slide, title)

        chart_type_name = cd.get("type", "column")
        chart_title = cd.get("title", "")
        categories = cd.get("categories", [])
        series_list = cd.get("series", [])

        # 图表放在卡片区域
        chart_top = Inches(1.5) if title else Inches(0.5)
        chart_height = Inches(5.6) if title else Inches(6.5)
        self._add_card(slide, Inches(0.5), chart_top,
                       Inches(12.333), chart_height)

        if categories and series_list:
            from pptx.chart.data import CategoryChartData

            xl_type = self.CHART_TYPES.get(chart_type_name, XL_CHART_TYPE.COLUMN_CLUSTERED)

            cdata = CategoryChartData()
            cdata.categories = categories
            for s in series_list:
                cdata.add_series(s.get("name", ""), s.get("values", []))

            chart_frame = slide.shapes.add_chart(
                xl_type,
                Inches(0.8), chart_top + Inches(0.3),
                Inches(11.733), chart_height - Inches(0.6),
                cdata,
            )
            chart = chart_frame.chart
            chart.has_legend = len(series_list) > 1
            if chart.has_legend:
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False

            for i, s in enumerate(chart.series):
                if i < len(t["chart_colors"]):
                    s.format.fill.solid()
                    s.format.fill.fore_color.rgb = t["chart_colors"][i]

            if chart_title:
                chart.has_title = True
                chart.chart_title.text_frame.paragraphs[0].text = chart_title
            else:
                chart.has_title = False

        self._add_notes(slide, data.get("notes", ""))
        self._apply_transition(slide, data.get("transition", ""))

    def _add_table(self, data: dict):
        """表格页：卡片内表格，更精致的表头/行样式"""
        t = self.theme
        slide = self._create_slide(t["bg_light"])

        title = data.get("title", "")
        td = data.get("table", {})
        headers = td.get("headers", [])
        rows = td.get("rows", [])

        if title:
            self._add_slide_title(slide, title)

        if headers and rows:
            num_rows = len(rows) + 1
            num_cols = len(headers)

            table_top = Inches(1.5) if title else Inches(0.5)

            # 先放卡片（表格上方留空间）
            card = self._add_card(slide, Inches(0.5), table_top,
                                  Inches(12.333), Inches(5.8))

            # 表格 - 盖在卡片上
            table_shape = slide.shapes.add_table(
                num_rows, num_cols,
                Inches(0.8), table_top + Inches(0.2),
                Inches(11.733), min(Inches(5.5), Inches(0.5 * num_rows)),
            )
            table = table_shape.table

            # 表头
            for ci in range(num_cols):
                cell = table.cell(0, ci)
                cell.text = str(headers[ci]) if ci < len(headers) else ""
                self._set_cell_bg(cell, t["primary"])
                for p in cell.text_frame.paragraphs:
                    self._set_paragraph_font(p, Pt(14), t["white"], bold=True)
                    p.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 数据行
            for ri, row in enumerate(rows):
                for ci, val in enumerate(row):
                    cell = table.cell(ri + 1, ci)
                    cell.text = str(val)
                    if ri % 2 == 0:
                        self._set_cell_bg(cell, RGBColor(0xF8, 0xF9, 0xFB))
                    for p in cell.text_frame.paragraphs:
                        self._set_paragraph_font(p, Pt(13), t["dark"])
                        p.alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        self._add_notes(slide, data.get("notes", ""))
        self._apply_transition(slide, data.get("transition", ""))

    def _add_section(self, data: dict):
        """章节页：深色背景 + 大字 + 装饰几何"""
        t = self.theme
        slide = self._create_slide(t["bg_dark"])

        sw = self.prs.slide_width

        # 顶部装饰条
        self._add_rect(slide, Inches(0), Inches(0), sw, Inches(0.04), t["accent"])

        # 左侧大字装饰（章节编号或首字母装饰）
        section_num = data.get("section_num", "")
        if section_num:
            tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.5),
                                          Inches(3), Inches(1.5))
            p = tb.text_frame.paragraphs[0]
            p.text = section_num
            self._set_paragraph_font(p, Pt(60), t["accent"], bold=True)
            p.alignment = PP_ALIGN.LEFT

        # 大标题
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(3.0),
                                      Inches(10.333), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        self._set_paragraph_font(p, Pt(40), t["white"], bold=True)
        p.alignment = PP_ALIGN.LEFT

        # 装饰横线
        self._add_rect(slide, Inches(1.5), Inches(5.3),
                       Inches(3), Inches(0.04), t["accent"])

        # 右下装饰方块
        self._add_rect(slide, sw - Inches(1.5), Inches(5.8),
                       Inches(0.8), Inches(0.8), t["accent"])

        self._apply_transition(slide, data.get("transition", ""))

    # ──────── 工具方法 ────────

    def _create_slide(self, bg_color):
        """创建空白幻灯片并设置背景色"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, bg_color)
        return slide

    def _set_bg(self, slide, color: RGBColor):
        """设置幻灯片背景色"""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_textbox(self, slide, left, top, width, height):
        """添加文本框"""
        return slide.shapes.add_textbox(left, top, width, height)

    def _add_rect(self, slide, left, top, width, height, color: RGBColor):
        """添加填充矩形（装饰线/色块）"""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_card(self, slide, left, top, width, height):
        """添加圆角白色卡片（带阴影）"""
        t = self.theme
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = t["bg_card"]
        shape.line.color.rgb = t["card_border"]
        shape.line.width = Pt(0.5)
        self._add_shadow(shape)
        return shape

    def _add_shadow(self, shape):
        """为形状添加柔和阴影（OXML）"""
        spPr = shape._element.spPr
        for child in list(spPr):
            if child.tag == qn('a:effectLst'):
                spPr.remove(child)

        effectLst = etree.SubElement(spPr, qn('a:effectLst'))
        outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
        outerShdw.set('blurRad', '38100')
        outerShdw.set('dist', '12700')
        outerShdw.set('dir', '2700000')
        outerShdw.set('algn', 'tl')
        outerShdw.set('rotWithShape', '0')

        srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
        srgbClr.set('val', '000000')
        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', '12000')

    def _add_slide_title(self, slide, title: str):
        """统一风格添加页面标题"""
        t = self.theme
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.35),
                                      Inches(11.733), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        self._set_paragraph_font(p, Pt(28), t["primary"], bold=True)

    def _add_column_text(self, slide, left, top, width, content, theme):
        """在卡片内添加栏文本"""
        tb = slide.shapes.add_textbox(left, top, width, Inches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)

        lines = content.split("\n") if content else []
        for li, line in enumerate(lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            line = line.strip()
            if not line:
                continue
            if line.startswith("## "):
                p.text = line[3:]
                self._set_paragraph_font(p, Pt(17), theme["accent"], bold=True)
                p.space_before = Pt(10)
            else:
                p.text = line
                self._set_paragraph_font(p, Pt(14), theme["dark"])
                p.space_after = Pt(4)
                p.line_spacing = Pt(22)

    def _add_notes(self, slide, notes: str):
        """添加演讲者备注"""
        if not notes:
            return
        try:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
        except Exception:
            pass

    def _apply_transition(self, slide, name: str):
        """应用幻灯片切换效果（OXML）"""
        if not name:
            return
        tag = self.TRANSITION_TAGS.get(name.lower())
        if not tag:
            return

        sld = slide._element
        trans_elem = sld.find(qn('p:transition'))
        if trans_elem is None:
            trans_elem = etree.SubElement(sld, qn('p:transition'))

        child_tag = qn(tag)
        existing = trans_elem.find(child_tag)
        if existing is None:
            etree.SubElement(trans_elem, child_tag)

        trans_elem.set('dur', '500')

    def _set_cell_bg(self, cell, color: RGBColor):
        """设置表格单元格背景色"""
        tcPr = cell._tc.get_or_add_tcPr()
        for child in list(tcPr):
            if child.tag == qn('a:solidFill'):
                tcPr.remove(child)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', str(color))

    def _set_paragraph_font(self, paragraph, size=None, color=None,
                            bold=None, name="Microsoft YaHei"):
        """设置段落字体（含东亚字体回退）"""
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        if bold is not None:
            run.font.bold = bold
        if size:
            run.font.size = size
        if color:
            run.font.color.rgb = color
        if name:
            run.font.name = name
            # 设置东亚字体
            rPr = run._r.get_or_add_rPr()
            for child in list(rPr):
                if child.tag == qn('a:ea'):
                    rPr.remove(child)
            ea = etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', name)


def create_presentation(
    title: str = "",
    slides: str = "",
    template_path: str = "",
    theme: str = "default",
    save_path: str = "",
) -> str:
    """创建 PowerPoint 演示文稿（完整功能版）

    参数:
        title: 演示标题（封面页用）
        slides: JSON 数组，每项格式见下方说明
        template_path: 模板文件路径 (.potx/.pptx)
        theme: 主题名 default|modern|minimal|dark|nature|warm
        save_path: 保存路径，默认桌面

    Slide JSON 格式:
        {
            "title": "页面标题",
            "content": "正文（\\n 换行，## 前缀=小标题）",
            "bullet": true,           // 项目符号模式（向后兼容）
            "layout": "content",      // content|bullet|two_column|image_right|image_left|section|chart|table
            "subtitle": "副标题",
            "section_num": "01",      // 章节编号（section 布局时用）
            "notes": "演讲者备注",
            "image": "C:/photo.png",  // 图片路径（图文布局时用）
            "chart": {                // 图表数据
                "type": "column",     // column|bar|line|pie|area
                "title": "图表标题",
                "categories": ["Q1","Q2"],
                "series": [{"name":"实际","values":[100,200]}]
            },
            "table": {                // 表格数据
                "headers": ["姓名","年龄"],
                "rows": [["张三","25"]]
            },
            "transition": "fade"      // none|fade|push|wipe|split|zoom
        }
    """
    try:
        builder = PresentationBuilder(template_path=template_path, theme=theme)
    except Exception as e:
        return f"错误：初始化演示文稿失败 - {e}"

    slide_list = []
    if slides:
        if isinstance(slides, list):
            # 直接传了 Python 列表，不再序列化
            slide_list = slides
        elif isinstance(slides, str):
            try:
                slide_list = json.loads(slides)
            except json.JSONDecodeError:
                try:
                    fixed = slides.replace('\n', '\\n').replace('\r', '')
                    slide_list = json.loads(fixed)
                except json.JSONDecodeError as e:
                    return f"错误：slides 参数 JSON 格式不正确 - {e}"
        else:
            return "错误：slides 参数必须是 JSON 字符串或 Python 列表"

    if not slide_list and not title:
        return "错误：请提供演示标题或至少一张幻灯片"

    # 封面
    if title:
        subtitle = ""
        if slide_list and isinstance(slide_list[0], dict):
            subtitle = slide_list[0].get("subtitle", "")
        builder.add_cover(title, subtitle)

    # 内容页
    for s in slide_list:
        try:
            builder.add_slide(s)
        except Exception as e:
            return f"错误：第 {slide_list.index(s) + 1} 页生成失败 - {e}"

    if not save_path:
        desktop = os.path.expanduser("~\\Desktop")
        safe_title = title.replace("/", "_").replace("\\", "_") if title else "演示文稿"
        save_path = os.path.join(desktop, f"{safe_title}.pptx")

    try:
        builder.save(save_path)
    except Exception as e:
        return f"错误：保存文件失败 - {e}"

    lines = [
        f"演示文稿已生成！",
        f"文件路径: {save_path}",
        f"幻灯片数: {builder.slide_count}",
        f"主题: {builder.theme['name']}",
    ]
    if template_path:
        lines.append(f"模板: {template_path}")
    return "\n".join(lines)
